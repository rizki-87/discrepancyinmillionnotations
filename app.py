from concurrent.futures import ThreadPoolExecutor
import streamlit as st
import tempfile
from pathlib import Path
from pptx import Presentation
import csv
import re
import logging

# Setup logging
logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(levelname)s - %(message)s')

# Password Protection
PREDEFINED_PASSWORD = "securepassword123"

def password_protection():
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if not st.session_state.authenticated:
        with st.form("password_form", clear_on_submit=True):
            password_input = st.text_input("Enter Password", type="password")
            submitted = st.form_submit_button("Submit")
            if submitted and password_input == PREDEFINED_PASSWORD:
                st.session_state.authenticated = True
                st.success("Access Granted! Please click 'Submit' again to proceed.")
            elif submitted:
                st.error("Incorrect Password")
        return False
    return True

# Million Notation Validation
def validate_million_notations(slide, slide_index):
    issues = []
    million_patterns = [r'\b\d+M\b', r'\b\d+\s?Million\b', r'\b\d+mn\b', r'\b\d+\sm\b']  # Patterns to match million notations
    notation_set = set()
    all_matches = []

    logging.debug(f"Slide {slide_index}: Checking shapes for million notations")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            logging.debug(f"Slide {slide_index}: Shape without text frame skipped")
            continue
        
        logging.debug(f"Slide {slide_index}: Text frame detected")
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for pattern in million_patterns:
                    matches = re.findall(pattern, run.text, re.IGNORECASE)
                    logging.debug(f"Slide {slide_index}: Found matches with pattern {pattern} - {matches}")
                    all_matches.extend(matches)
                    for match in matches:
                        notation_set.add(pattern)

    logging.debug(f"Slide {slide_index}: Notation set - {notation_set}")
    if len(notation_set) > 1:
        for match in all_matches:
            logging.debug(f"Slide {slide_index}: Inconsistent million notation detected - {match}")
            issues.append({
                'slide': slide_index,
                'issue': 'Inconsistent Million Notations',
                'text': match,
                'details': f'Found inconsistent million notations: {list(notation_set)}'
            })

    return issues

# Save results to CSV
def save_to_csv(issues, output_csv):
    with open(output_csv, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.DictWriter(file, fieldnames=['slide', 'issue', 'text', 'corrected', 'details'])
        writer.writeheader()
        writer.writerows(issues)

def main():
    if not password_protection():
        return

    st.title("Million Notations Validator")
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])

    if uploaded_file:
        with tempfile.TemporaryDirectory() as tmpdir:
            temp_ppt_path = Path(tmpdir) / "uploaded_ppt.pptx"
            with open(temp_ppt_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            presentation = Presentation(temp_ppt_path)
            total_slides = len(presentation.slides)
            logging.debug(f"Total slides: {total_slides}")

            # Run Validation
            if st.button("Run Million Notations Validation"):
                progress_bar = st.progress(0)
                progress_text = st.empty()
                issues = []

                # Manual Iteration
                for slide_index in range(total_slides):
                    slide = presentation.slides[slide_index]
                    logging.debug(f"Validating slide {slide_index + 1}")
                    slide_issues = validate_million_notations(slide, slide_index + 1)
                    issues.extend(slide_issues)
                    progress_percent = int((slide_index + 1) / total_slides * 100)
                    progress_text.text(f"Progress: {progress_percent}%")
                    progress_bar.progress(progress_percent / 100)

                # Save Results
                csv_output_path = Path(tmpdir) / "million_notations_validation_report.csv"
                save_to_csv(issues, csv_output_path)

                # Store results in session state
                st.session_state['csv_output'] = csv_output_path.read_bytes()

                st.success("Million notations validation completed!")

            # Display Download Buttons
            if 'csv_output' in st.session_state:
                st.download_button("Download Validation Report (CSV)", st.session_state['csv_output'],
                                   file_name="million_notations_validation_report.csv")

if __name__ == "__main__":
    main()
